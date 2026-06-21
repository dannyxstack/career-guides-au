"""从职业数据生成有图有表的 PPTX，配 Pexels 在线职业图片。

用法：
    python -m video_pipeline.ppt_builder <occupation_id> [--locale zh-CN] [--no-images]

输出：video_pipeline/out/ppt/<slug>.pptx
"""

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from video_pipeline import charts, config, content_store, images, scriptwriter  # noqa: E402
from video_pipeline.data_source import get_occupation  # noqa: E402

# ---- 配色（与视频主题一致）----
BG = RGBColor(0x0B, 0x11, 0x20)
CARD = RGBColor(0x15, 0x1C, 0x2C)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xE6, 0xEB, 0xF5)
MUTED = RGBColor(0x8B, 0x97, 0xB0)
LINE = RGBColor(0x26, 0x30, 0x4A)

CJK_FONT = "Microsoft YaHei"
SW, SH = Inches(13.333), Inches(7.5)  # 16:9

# 评分维度 → 中文显示名（让评分页能看懂评的是什么）
DIM_ZH = {
    "learning_difficulty": "学习难度",
    "learning_duration": "学习时长",
    "certification_difficulty": "认证难度",
    "job_demand": "岗位需求",
    "competition": "竞争程度",
    "work_intensity": "工作强度",
    "income_level": "收入水平",
    "future_prospect": "未来前景",
    "ai_risk": "AI替代风险",
    "pr_friendliness": "移民友好度",
    "pr_difficulty": "移民难度",
}


def _font(run, size, color=WHITE, bold=False, name=CJK_FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    # 同时设置东亚字体，避免中文变方块
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {"typeface": name})
    rPr.append(ea)


def _slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def _text(slide, x, y, w, h, lines, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: [(text, size, color, bold), ...]"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = txt
        _font(run, size, color, bold)
    return tb


def _title_bar(slide, title):
    """页面左上角的绿色标题。"""
    bar = slide.shapes.add_shape(1, Inches(0.7), Inches(0.55), Inches(0.12), Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()
    _text(slide, Inches(0.95), Inches(0.5), Inches(11), Inches(0.7),
          [(title, 30, WHITE, True)])


def _pic_cover(slide, path, x, y, w, h):
    """放图片到指定框（填满，轻微裁切由 PPT 处理）。"""
    if not path or not Path(path).exists():
        return
    slide.shapes.add_picture(path, x, y, width=w, height=h)


def _slug(s):
    return re.sub(r"-+", "-", re.sub(r"[^a-z0-9]+", "-", (s or "occ").lower())).strip("-")


# ---------- 各页 ----------

def _slide_title(prs, occ, imgs, brief=None):
    s = _slide(prs)
    # 右侧实拍图占 ~45%，避免整页深色文字
    if imgs:
        _pic_cover(s, imgs[0], Inches(7.4), 0, Inches(5.93), SH)
    hook = (brief or {}).get("hook_title") or occ["name_zh"] or occ["name_en"] or ""
    lines = [(hook, 40, WHITE, True),
             (f"{occ['name_zh'] or ''} · {occ['name_en'] or ''}", 20, MUTED, False)]
    # 首屏三个关键数字/卖点
    for stat in (brief or {}).get("key_stats", [])[:3]:
        lines.append((f"▸ {stat}", 18, GREEN, True))
    _text(s, Inches(0.8), Inches(1.7), Inches(6.3), Inches(4.6), lines)


def _slide_daily(prs, occ, imgs, brief):
    """这个职业每天做什么（决策视频式，紧跟钩子）。"""
    s = _slide(prs)
    _title_bar(s, "这个职业每天做什么")
    tasks = (brief or {}).get("daily_tasks") or []
    lines = [(f"• {t}", 20, WHITE, False) for t in tasks] or \
            [(occ["summary"] or "（暂无数据）", 20, WHITE, False)]
    _text(s, Inches(0.95), Inches(1.7), Inches(6.3), Inches(5), lines)
    img = imgs[1] if len(imgs) > 1 else (imgs[0] if imgs else None)
    if img:
        _pic_cover(s, img, Inches(7.7), Inches(1.7), Inches(5), Inches(3.3))


def _slide_suitability(prs, occ):
    """适合谁 / 不适合谁（双栏）。数据来自 DB suitability_fit/unfit。"""
    s = _slide(prs)
    _title_bar(s, "适合谁 / 不适合谁")
    fit = occ.get("suitability_fit") or []
    unfit = occ.get("suitability_unfit") or []
    _text(s, Inches(0.95), Inches(1.6), Inches(5.8), Inches(0.6),
          [("✓ 适合", 22, GREEN, True)])
    _text(s, Inches(0.95), Inches(2.2), Inches(5.8), Inches(4.5),
          [(f"• {x}", 18, WHITE, False) for x in fit] or [("—", 18, MUTED, False)])
    _text(s, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.6),
          [("✗ 不适合", 22, AMBER, True)])
    _text(s, Inches(7.0), Inches(2.2), Inches(5.6), Inches(4.5),
          [(f"• {x}", 18, MUTED, False) for x in unfit] or [("—", 18, MUTED, False)])


def _slide_action(prs, occ, brief):
    """30 天入门行动清单。"""
    s = _slide(prs)
    _title_bar(s, "30 天入门行动清单")
    steps = (brief or {}).get("action_plan_30d") or []
    lines = [(f"{i}. {st}", 20, WHITE, False) for i, st in enumerate(steps, 1)] or \
            [("（暂无行动清单，请先生成 PPT 文案 brief）", 18, MUTED, False)]
    _text(s, Inches(0.95), Inches(1.7), Inches(11.4), Inches(5), lines)


def _table(slide, rows, cols, x, y, w, h, header):
    tbl = slide.shapes.add_table(rows, cols, x, y, w, h).table
    for j, head in enumerate(header):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = GREEN
        c.text_frame.word_wrap = True
        r = c.text_frame.paragraphs[0].add_run(); r.text = head
        _font(r, 16, BG, True)
    return tbl


def _fill_cell(cell, text, color=WHITE, bg=CARD, size=15, bold=False):
    cell.fill.solid(); cell.fill.fore_color.rgb = bg
    cell.text_frame.word_wrap = True
    r = cell.text_frame.paragraphs[0].add_run(); r.text = text
    _font(r, size, color, bold)


def _slide_salary(prs, occ):
    s = _slide(prs)
    _title_bar(s, "薪资范围（AUD / 年）")
    sal = occ["salaries"]
    tbl = _table(s, len(sal) + 1, 2, Inches(0.95), Inches(1.7),
                 Inches(11.4), Inches(0.6 + 0.7 * len(sal)), ["经验阶段", "年薪范围"])
    for i, row in enumerate(sal, start=1):
        bg = CARD if i % 2 else BG
        _fill_cell(tbl.cell(i, 0), str(row["label"]), WHITE, bg)
        rng = f"${int(row['min']):,} ~ ${int(row['max']):,}" if row["min"] and row["max"] else "—"
        _fill_cell(tbl.cell(i, 1), rng, AMBER, bg, bold=True)


def _slide_ratings(prs, occ):
    """全部维度用放射（雷达）图呈现，标题展示综合评分（如 7.8/10）。"""
    s = _slide(prs)
    score = charts.overall_score(occ)
    title = "职业评分" + (f" · 综合 {score}/10" if score is not None else "")
    _title_bar(s, title)
    slug = _slug(occ["name_en"] or str(occ["occupation_id"]))
    png = charts.rating_radar_png(
        occ, config.OUT_DIR / "img" / f"{slug}_radar.png", DIM_ZH)
    if png and Path(png).exists():
        # 雷达图近似正方形，水平居中
        s.shapes.add_picture(png, Inches(4.0), Inches(1.35), height=Inches(5.7))
    else:
        _text(s, Inches(0.95), Inches(1.7), Inches(11.4), Inches(5),
              [("（暂无评分数据）", 18, MUTED, False)])


def _slide_visa(prs, occ, brief=None):
    s = _slide(prs)
    _title_bar(s, "移民路径 / ANZSCO 风险提示")
    lines = []
    for v in occ["visa_pathways"]:
        lines.append((f"● {v['subclass']}  {v['name'] or ''}", 20, GREEN, True))
        if v.get("desc"):
            lines.append((f"    {v['desc']}", 16, MUTED, False))
    if not lines:
        lines = [("（暂无签证数据）", 18, MUTED, False)]
    _text(s, Inches(0.95), Inches(1.7), Inches(11.4), Inches(3.6), lines)
    # 风险提示：签证须按具体职责匹配 ANZSCO，以官方与评估机构为准
    note = (brief or {}).get("anzsco_note") or (
        "签证路径需按具体职责匹配对应 ANZSCO 职业，并以 Department of Home Affairs "
        "最新职业清单及相关评估机构（如 ACS）评估结果为准。")
    _text(s, Inches(0.95), Inches(5.7), Inches(11.4), Inches(1.4),
          [("⚠ " + note, 15, AMBER, False)])


def _slide_growth(prs, occ, imgs, brief=None):
    s = _slide(prs)
    _title_bar(s, "增长方向 / 搜索热词")
    # 优先用 LLM brief 里贴近搜索流量的关键词，回退到 DB 的 growth_areas
    kws = (brief or {}).get("growth_keywords") or occ["growth_areas"]
    lines = [(f"▸ {g}", 22, WHITE, False) for g in kws] or \
            [("（暂无数据）", 18, MUTED, False)]
    _text(s, Inches(0.95), Inches(1.7), Inches(6.3), Inches(5), lines)
    img = imgs[2] if len(imgs) > 2 else (imgs[-1] if imgs else None)
    if img:
        _pic_cover(s, img, Inches(7.7), Inches(1.7), Inches(5), Inches(3.3))




def _fmt_cost(e):
    try:
        lo, hi = float(e.get("cost_min") or 0), float(e.get("cost_max") or 0)
    except (TypeError, ValueError):
        lo = hi = 0
    if hi > 0:
        return f"${int(lo):,}~${int(hi):,}"
    return e.get("cost_note") or "—"


def _slide_education(prs, occ):
    s = _slide(prs)
    _title_bar(s, "教育路径")
    edu = occ["education"]
    if not edu:
        _text(s, Inches(0.95), Inches(1.7), Inches(11), Inches(1),
              [("（暂无教育路径数据）", 20, MUTED, False)])
        return
    tbl = _table(s, len(edu) + 1, 3, Inches(0.95), Inches(1.7),
                 Inches(11.4), Inches(0.6 + 0.7 * len(edu)), ["阶段", "时长", "费用 (AUD)"])
    for i, e in enumerate(edu, start=1):
        bg = CARD if i % 2 else BG
        _fill_cell(tbl.cell(i, 0), str(e["stage"]), WHITE, bg)
        _fill_cell(tbl.cell(i, 1), str(e["duration"] or "—"), WHITE, bg)
        _fill_cell(tbl.cell(i, 2), _fmt_cost(e), AMBER, bg, bold=True)


def _slide_qualifications(prs, occ):
    s = _slide(prs)
    _title_bar(s, "从业资质 / 证书")
    lines = []
    for q in occ["qualifications"]:
        tag = "（必备）" if q["mandatory"] else "（可选）"
        lines.append((f"● {q['name']}  {tag}", 20,
                      GREEN if q["mandatory"] else MUTED, True))
        if q.get("issuer"):
            lines.append((f"    颁发：{q['issuer']}", 15, MUTED, False))
    if not lines:
        lines = [("（暂无资质数据）", 18, MUTED, False)]
    _text(s, Inches(0.95), Inches(1.7), Inches(11.4), Inches(5), lines)


def _slide_income(prs, occ, slug):
    s = _slide(prs)
    _text(s, Inches(0.8), Inches(0.45), Inches(12), Inches(0.9),
          [("收入趋势：现在入行，正赶上上升期", 28, GREEN, True)])
    png = charts.salary_trend_png(occ, config.OUT_DIR / "img" / f"{slug}_salary.png")
    if Path(png).exists():
        s.shapes.add_picture(png, Inches(0.55), Inches(1.5), width=Inches(12.2))


def _slide_demand(prs, occ, slug):
    s = _slide(prs)
    _text(s, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
          [("岗位需求估算：2025 后仍偏强，但初级竞争加剧", 26, GREEN, True)])
    png = charts.demand_trend_png(occ, config.OUT_DIR / "img" / f"{slug}_demand.png")
    if Path(png).exists():
        s.shapes.add_picture(png, Inches(0.55), Inches(1.35), width=Inches(12.2))
    _text(s, Inches(0.8), Inches(6.95), Inches(12), Inches(0.5),
          [("估算逻辑：以当前岗位量、薪资区间、行业数字化趋势生成，不代表官方预测。", 12, MUTED, False)])


def _slide_cta(prs, occ):
    s = _slide(prs)
    _text(s, Inches(1), Inches(2.6), Inches(11.3), Inches(2.5),
          [("想了解更多澳洲紧缺职业？", 36, WHITE, True),
           ("关注我们，每天一个职业深度解析。", 24, GREEN, False),
           ("数据来源：Jobs and Skills Australia · Seek · Home Affairs", 14, MUTED, False)],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _get_brief(occupation_id: int, locale: str, auto: bool = True) -> dict | None:
    """读取已入库的 PPT 文案 brief；没有且 auto=True 时用 LLM 生成（失败则降级返回 None）。"""
    row = content_store.get_content_by_key(occupation_id, locale, "ppt")
    if row and row.get("outline"):
        return row["outline"]
    if not auto:
        return None
    try:
        return scriptwriter.generate_ppt_brief(occupation_id, locale)
    except Exception as e:  # noqa: BLE001 — brief 是增强项，失败不应中断出 PPT
        print(f"[ppt] 生成文案 brief 失败，降级为无新内容版本：{e}")
        return None


def build_ppt(occupation_id: int, locale: str = "zh-CN", with_images: bool = True,
              out_path: str | None = None, brief: dict | None = None) -> str:
    occ = get_occupation(occupation_id, locale=locale)
    if brief is None:
        brief = _get_brief(occupation_id, locale)
    imgs = []
    if with_images:
        q = occ["name_en"] or occ["name_zh"] or "professional worker"
        imgs = images.search_images(q, count=3)
        if len(imgs) < 3:
            imgs += images.search_images(f"{q} at work", count=3 - len(imgs))
        print(f"[images] 获取 {len(imgs)} 张图片")

    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    slug = _slug(occ["name_en"] or str(occupation_id))

    # 决策视频式顺序：钩子 → 每天做什么 → 薪资与岗位量 → 评分 → 适合谁 →
    # 技能路线 → 增长热词 → 移民/ANZSCO 风险 → 30天行动清单 → CTA
    _slide_title(prs, occ, imgs, brief)
    _slide_daily(prs, occ, imgs, brief)
    _slide_income(prs, occ, slug)         # 收入趋势图
    _slide_demand(prs, occ, slug)         # 需求估算图
    _slide_salary(prs, occ)               # 薪资表
    _slide_ratings(prs, occ)              # 5 维评分速览
    _slide_suitability(prs, occ)          # 适合谁 / 不适合谁
    _slide_education(prs, occ)            # 技能路线：教育路径
    _slide_qualifications(prs, occ)       # 技能路线：从业资质 / 证书
    _slide_growth(prs, occ, imgs, brief)  # 增长方向 / 搜索热词
    _slide_visa(prs, occ, brief)          # 移民路径 / ANZSCO 风险
    _slide_action(prs, occ, brief)        # 30 天入门行动清单
    _slide_cta(prs, occ)

    if out_path:
        final = Path(out_path)
    else:
        out_dir = config.OUT_DIR / "ppt"
        out_dir.mkdir(parents=True, exist_ok=True)
        final = out_dir / f"{_slug(occ['name_en'] or str(occupation_id))}.pptx"
    final.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(final))
    print(f"[ppt] {final}  ({len(prs.slides)} slides)")
    return str(final)


if __name__ == "__main__":
    import argparse

    ap = argparse.ArgumentParser()
    ap.add_argument("occupation_id", type=int)
    ap.add_argument("--locale", default="zh-CN")
    ap.add_argument("--no-images", action="store_true")
    ap.add_argument("--out", default=None, help="自定义输出 pptx 路径")
    args = ap.parse_args()
    path = build_ppt(args.occupation_id, args.locale,
                     with_images=not args.no_images, out_path=args.out)
    print(f"[OK] {path}")
